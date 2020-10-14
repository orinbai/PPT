from pptx import Presentation
from pptx.chart.data import CategoryChartData

prs = Presentation("th.pptx")
core = prs.core_properties
# chartData = CategoryChartData()
cateData = ["A3", "A4", "A5"]
seriesData = [
    ["2010", (20.1, 22.1, 25.2)],
    ["2011", (15.2, 12.1, 13.5)],
    ["2012", (19.1, 18.1, 17.2)],
    ["2010", (23.1, 25.1, 24.2)]
]
# chartData.add_series("2010", (20.1, 22.1, 25.2))
# chartData.add_series("2011", (15.2, 12.1, 13.5))
# chartData.add_series("2012", (19.1, 18.1, 17.2))
# chartData.add_series("2010", (23.1, 25.1, 24.2))


def reNewChart(obj, cateData, seData):
    chartData = CategoryChartData()
    chartData.categories = cateData
    for sdata in seData:
        chartData.add_series(sdata[0], sdata[1])
    obj.replace_data(chartData)


print(len(prs.slides[0].shapes))
for n, shape in enumerate(prs.slides[0].shapes):
    print(shape)
    print(n, shape.shape_id, shape.has_chart)
    if shape.has_chart:
        old_id = shape.shape_id
        print("="*20)
        print(prs.slides[0].shapes[n].chart)
        print("="*20)
        reNewChart(prs.slides[0].shapes[n].chart, cateData, seriesData)
prs.save("test.1.pptx")
